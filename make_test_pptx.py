# -*- coding: utf-8 -*-
"""Generate a realistic departmental weekly-report PPTX for testing the tool."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import io, struct, zlib, os

def make_png(w, h, rgb):
    """Build a tiny solid-color PNG in pure python (no PIL needed)."""
    def chunk(typ, data):
        c = typ + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    sig = b'\x89PNG\r\n\x1a\n'
    ihdr = struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0)  # 8-bit RGB
    raw = b''
    for y in range(h):
        raw += b'\x00' + bytes(rgb) * w
    idat = zlib.compress(raw, 9)
    png = sig + chunk(b'IHDR', ihdr) + chunk(b'IDAT', idat) + chunk(b'IEND', b'')
    return png

# Each entry: one project row in the weekly report.
ROWS = [
    {
        "Reporter": "David",
        "Project Name": "Sensor Firmware v2.1",
        "Current Job and Issue": "Completed I2C driver refactor; intermittent ACK timeout still seen on cold boot, suspect pull-up timing.",
        "Risk": "High",
        "Due Date": "2026-06-20",
        "Owner": "Eric/Isha",
        "Next Week Job And Plan": "Add oscilloscope capture on SDA line; finalize timing fix and regression test.",
        "img": (200, 120, (60, 120, 200)),
    },
    {
        "Reporter": "David",
        "Project Name": "Power Rail Validation",
        "Current Job and Issue": "Bench measurement of 3V3 rail ripple under load; waiting on new load board from vendor.",
        "Risk": "Medium",
        "Due Date": "2026-06-27",
        "Owner": "Rick, Jin, John",
        "Next Week Job And Plan": "Repeat ripple test once load board arrives; document margins.",
        "img": None,
    },
    {
        "Reporter": "Mei",
        "Project Name": "Cloud Telemetry Pipeline",
        "Current Job and Issue": "Debugging dropped MQTT packets at high throughput; added retry queue, verifying delivery.",
        "Risk": "Low",
        "Due Date": "2026-07-04",
        "Owner": "Sam / Adrian",
        "Next Week Job And Plan": "Load test at 10k msg/s; add Grafana dashboard for queue depth.",
        "img": (220, 100, (200, 80, 80)),
    },
    {
        "Reporter": "Mei",
        "Project Name": "Mobile App Login",
        "Current Job and Issue": "OAuth refresh token rotation implemented; document study on PKCE flow ongoing.",
        "Risk": "Low",
        "Due Date": "2026-06-30",
        "Owner": "Ravi",
        "Next Week Job And Plan": "Verification of token rotation across 3 devices; write integration tests.",
        "img": None,
    },
    {
        "Reporter": "Ken",
        "Project Name": "Thermal Chamber Automation",
        "Current Job and Issue": "Blocked by vendor: chamber SDK license not yet delivered. Prepared test scripts in advance.",
        "Risk": "High",
        "Due Date": "2026-07-11",
        "Owner": "Raveendra",
        "Next Week Job And Plan": "Once license arrives, dry-run automated soak test; otherwise escalate to procurement.",
        "img": (180, 140, (90, 170, 90)),
    },
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for r in ROWS:
    slide = prs.slides.add_slide(blank)
    # Title = project
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8))
    tf = tb.text_frame
    tf.text = r["Project Name"]
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True

    # A 2-column table: field name | value
    fields = ["Reporter", "Project Name", "Current Job and Issue", "Risk",
              "Due Date", "Owner", "Next Week Job And Plan"]
    rows = len(fields)
    tbl_shape = slide.shapes.add_table(rows, 2, Inches(0.4), Inches(1.1),
                                       Inches(8.5), Inches(5.5))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(5.9)
    for i, f in enumerate(fields):
        table.cell(i, 0).text = f
        table.cell(i, 1).text = str(r[f])
        for c in (0, 1):
            for p in table.cell(i, c).text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
        table.cell(i, 0).text_frame.paragraphs[0].runs[0].font.bold = True

    # Optional image on the right side
    if r["img"]:
        w, h, color = r["img"]
        png = make_png(w, h, color)
        slide.shapes.add_picture(io.BytesIO(png), Inches(9.2), Inches(1.3),
                                 width=Inches(3.6))

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "test.pptx")
prs.save(out)
print("wrote", out, os.path.getsize(out), "bytes,", len(prs.slides._sldIdLst), "slides")
